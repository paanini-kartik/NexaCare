"""
Build NexaCare_Pitch_5.pptx — a condensed 5-slide version of the original deck.
Matches the exact visual language: #0f172a bg, #3b82f6 blue, #94a3b8 slate, white.
Card pattern: #1e293b fill + 4px left blue accent strip.
Stat card: #1e293b card + blue left strip + giant blue number + slate description.
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

# ---------------------------------------------------------------------------
# Colour constants (all as RGBColor)
# ---------------------------------------------------------------------------
NAVY   = RGBColor(0x0F, 0x17, 0x2A)   # #0f172a
SLATE1 = RGBColor(0x1E, 0x29, 0x3B)   # #1e293b (card bg)
BLUE   = RGBColor(0x3B, 0x82, 0xF6)   # #3b82f6
SLATE2 = RGBColor(0x94, 0xA3, 0xB8)   # #94a3b8 (secondary text)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
SLATE3 = RGBColor(0x0F, 0x23, 0x3F)   # slightly lighter navy for B2B tint

# Slide dimensions: 13.33" × 7.5" = 12192000 × 6858000 EMU
W = 12192000
H = 6858000

# ---------------------------------------------------------------------------
# Helper: create a new blank presentation (widescreen)
# ---------------------------------------------------------------------------
def new_prs():
    prs = Presentation()
    prs.slide_width  = Emu(W)
    prs.slide_height = Emu(H)
    return prs

# ---------------------------------------------------------------------------
# Low-level shape builders (work directly with python-pptx slide.shapes)
# ---------------------------------------------------------------------------
def add_rect(slide, x, y, w, h, fill_rgb, no_line=True):
    """Add a filled rectangle, return the shape."""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Emu(x), Emu(y), Emu(w), Emu(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if no_line:
        shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h):
    """Add an empty textbox, return the shape."""
    from pptx.util import Emu
    return slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))


def set_para(tf, text, size_pt, bold=False, italic=False, color=WHITE, align=PP_ALIGN.LEFT, space_before=0):
    """Add a paragraph to a text frame."""
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def clear_tf(tf):
    """Remove all paragraphs from a text frame."""
    from pptx.oxml.ns import qn
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    # Clear first paragraph
    tf.paragraphs[0].clear()


# ---------------------------------------------------------------------------
# Reusable component builders
# ---------------------------------------------------------------------------

def bg(slide):
    """Full-slide navy background."""
    add_rect(slide, 0, 0, W, H, NAVY)


def section_pill(slide, label, x=365760, y=274320, w=2194560, h=320040):
    """Small slate pill with uppercase section label (e.g. THE PROBLEM)."""
    add_rect(slide, x, y, w, h, SLATE1)
    tb = add_textbox(slide, x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = False
    clear_tf(tf)
    set_para(tf, label, 9, bold=True, color=SLATE2, align=PP_ALIGN.CENTER)


def slide_title(slide, text, x=548640, y=822960, w=10972800, h=640080):
    """Large white heading."""
    tb = add_textbox(slide, x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    clear_tf(tf)
    set_para(tf, text, 38, bold=True, color=WHITE)


def blue_underline(slide, x=548640, y=1508760, w=2743200, h=36576):
    """Thin electric-blue underline beneath title."""
    add_rect(slide, x, y, w, h, BLUE)


def stat_card(slide, x, y, w, h, number_text, desc_text):
    """
    Stat card: #1e293b background + 4px blue left accent strip.
    Giant blue number at top, slate description below.
    """
    STRIP = 54864   # ~6px
    # Card body
    add_rect(slide, x, y, w, h, SLATE1)
    # Blue left strip
    add_rect(slide, x, y, STRIP, h, BLUE)

    # Number textbox — occupies top ~40% of card
    num_h = int(h * 0.42)
    tb_num = add_textbox(slide, x + STRIP + 91440, y + 228600, w - STRIP - 182880, num_h)
    tf = tb_num.text_frame
    tf.word_wrap = False
    clear_tf(tf)
    set_para(tf, number_text, 52, bold=True, color=BLUE)

    # Description textbox — occupies bottom ~45%
    desc_y = y + num_h + 91440
    desc_h = h - num_h - 91440
    tb_desc = add_textbox(slide, x + STRIP + 91440, desc_y, w - STRIP - 182880, desc_h)
    tf2 = tb_desc.text_frame
    tf2.word_wrap = True
    clear_tf(tf2)
    set_para(tf2, desc_text, 14, bold=False, color=SLATE2)


def feature_card(slide, x, y, w, h, icon, title, body):
    """
    Feature card: #1e293b bg + 4px blue left accent.
    Icon+title as white bold, body as slate.
    """
    STRIP = 54864
    add_rect(slide, x, y, w, h, SLATE1)
    add_rect(slide, x, y, STRIP, h, BLUE)

    inner_x = x + STRIP + 137160
    inner_w = w - STRIP - 228600
    inner_y = y + 182880

    tb = add_textbox(slide, inner_x, inner_y, inner_w, h - 365760)
    tf = tb.text_frame
    tf.word_wrap = True
    clear_tf(tf)
    set_para(tf, f"{icon}  {title}", 14, bold=True, color=WHITE)
    set_para(tf, body, 11, bold=False, color=SLATE2, space_before=4)


def blue_pill_badge(slide, x, y, text, w=None):
    """Blue-bordered pill badge (for role labels or hackathon badge)."""
    from pptx.util import Pt
    est_w = w if w else (len(text) * 120000 + 200000)
    shape = slide.shapes.add_shape(1, Emu(x), Emu(y), Emu(est_w), Emu(228600))
    shape.fill.background()
    shape.line.color.rgb = BLUE
    shape.line.width = Emu(19050)

    tf = shape.text_frame
    tf.word_wrap = False
    clear_tf(tf)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = WHITE
    return shape


def caption(slide, text, x, y, w, h, italic=True, align=PP_ALIGN.CENTER, size=13):
    """Italic slate caption."""
    tb = add_textbox(slide, x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    clear_tf(tf)
    set_para(tf, text, size, bold=False, italic=italic, color=SLATE2, align=align)


# ---------------------------------------------------------------------------
# SLIDE 1 — HOOK
# ---------------------------------------------------------------------------
def build_slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Background
    bg(slide)

    # Top-right badge: "Hack the Globe 2026"
    badge_w = 2194560
    badge_x = W - badge_w - 274320
    badge_y = 274320
    badge_h = 365760
    add_rect(slide, badge_x, badge_y, badge_w, badge_h, BLUE)
    tb = add_textbox(slide, badge_x, badge_y, badge_w, badge_h)
    tf = tb.text_frame
    tf.word_wrap = False
    clear_tf(tf)
    set_para(tf, "Hack the Globe 2026", 11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Center blue circle (logomark)
    circle_size = 1645920
    circle_x = (W - circle_size) // 2
    circle_y = 1280160
    circ = slide.shapes.add_shape(
        9,  # OVAL
        Emu(circle_x), Emu(circle_y), Emu(circle_size), Emu(circle_size)
    )
    circ.fill.solid()
    circ.fill.fore_color.rgb = BLUE
    circ.line.fill.background()

    # "N" in the circle
    tb_n = add_textbox(slide, circle_x, circle_y, circle_size, circle_size)
    tf_n = tb_n.text_frame
    tf_n.word_wrap = False
    clear_tf(tf_n)
    p = tf_n.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "N"
    run.font.size = Pt(72)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # Title: NexaCare
    title_y = circle_y + circle_size + 228600
    tb_title = add_textbox(slide, 0, title_y, W, 731520)
    tf_t = tb_title.text_frame
    tf_t.word_wrap = False
    clear_tf(tf_t)
    set_para(tf_t, "NexaCare", 52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle (italic blue)
    sub_y = title_y + 640080
    tb_sub = add_textbox(slide, 0, sub_y, W, 365760)
    tf_s = tb_sub.text_frame
    tf_s.word_wrap = False
    clear_tf(tf_s)
    set_para(tf_s, "\u201cYour benefits exist. Start using them.\u201d", 18, bold=False, italic=True, color=BLUE, align=PP_ALIGN.CENTER)

    # Bottom row: 4 role badges
    roles = ["Frontend", "Backend", "Maps", "AI"]
    pill_w = 1645920
    pill_h = 274320
    total_pill_w = len(roles) * pill_w + (len(roles) - 1) * 182880
    start_x = (W - total_pill_w) // 2
    pill_y = H - pill_h - 457200
    for i, role in enumerate(roles):
        px = start_x + i * (pill_w + 182880)
        blue_pill_badge(slide, px, pill_y, role, w=pill_w)


# ---------------------------------------------------------------------------
# SLIDE 2 — PROBLEM
# ---------------------------------------------------------------------------
def build_slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    section_pill(slide, "THE PROBLEM")
    slide_title(slide, "The friction is the problem")
    blue_underline(slide)

    # Three stat cards side-by-side
    card_w = 3566880
    card_h = 1554480
    gap = 274320
    total_w = 3 * card_w + 2 * gap
    start_x = (W - total_w) // 2
    card_y = 1645920

    stats = [
        ("42%",  "of adults skipped a\npreventive checkup last year"),
        ("67%",  "of employees don't\nfully use employer benefits"),
        ("$88B", "spent annually on\npreventable conditions"),
    ]
    for i, (num, desc) in enumerate(stats):
        cx = start_x + i * (card_w + gap)
        stat_card(slide, cx, card_y, card_w, card_h, num, desc)

    # Chaos map — 5 nodes around a central "You" node
    # Using simple rectangles / ellipses and connector lines
    node_y_base = card_y + card_h + 274320
    node_area_h = H - node_y_base - 320040
    center_x = W // 2
    center_y = node_y_base + node_area_h // 2

    node_w = 1097280
    node_h = 320040

    # Central "You" node
    you_x = center_x - node_w // 2
    you_y = center_y - node_h // 2
    you_rect = slide.shapes.add_shape(1, Emu(you_x), Emu(you_y), Emu(node_w), Emu(node_h))
    you_rect.fill.solid()
    you_rect.fill.fore_color.rgb = BLUE
    you_rect.line.fill.background()
    tb_you = add_textbox(slide, you_x, you_y, node_w, node_h)
    tf_you = tb_you.text_frame
    tf_you.word_wrap = False
    clear_tf(tf_you)
    set_para(tf_you, "You", 13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Five surrounding nodes
    import math
    outer_nodes = [
        ("Dental Portal \U0001F512", -110),
        ("Vision Portal \U0001F512", -40),
        ("Physio Portal \U0001F512", 30),
        ("Benefits PDF \u2753", 100),
        ("Annual Reset \u23F1", 170),
    ]
    radius = 1828800
    for label, angle_deg in outer_nodes:
        angle = math.radians(angle_deg)
        nx = center_x + int(radius * math.cos(angle)) - node_w // 2
        ny = center_y + int(radius * math.sin(angle)) - node_h // 2

        n_rect = slide.shapes.add_shape(1, Emu(nx), Emu(ny), Emu(node_w), Emu(node_h))
        n_rect.fill.solid()
        n_rect.fill.fore_color.rgb = SLATE1
        n_rect.line.color.rgb = BLUE
        n_rect.line.width = Emu(19050)

        tb_n = add_textbox(slide, nx, ny, node_w, node_h)
        tf_n = tb_n.text_frame
        tf_n.word_wrap = False
        clear_tf(tf_n)
        set_para(tf_n, label, 10, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

    # Caption at bottom
    cap_y = H - 400000
    caption(slide,
            "Three portals. A document you never opened. A reset date you always miss.",
            548640, cap_y, W - 1097280, 320040,
            italic=True, align=PP_ALIGN.CENTER, size=13)


# ---------------------------------------------------------------------------
# SLIDE 3 — SOLUTION + HOW IT WORKS
# ---------------------------------------------------------------------------
def build_slide3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    section_pill(slide, "THE SOLUTION")
    slide_title(slide, "One platform. Set up in two minutes.")
    blue_underline(slide, w=4480560)

    # 2×2 feature card grid
    col_gap = 228600
    row_gap = 182880
    top_y   = 1737360
    grid_w  = W - 1097280
    start_x = 548640
    card_w  = (grid_w - col_gap) // 2
    card_h  = 1188720

    features = [
        ("\U0001F4C5", "Smart Scheduling",  "AI flags overdue checkups. Book in seconds."),
        ("\U0001F4B0", "Benefits Tracker",  "Dental, vision, physio \u2014 used vs. remaining."),
        ("\U0001F5FA",  "Health Compass",    "Nearby clinics filtered by your actual coverage."),
        ("\U0001F916", "AI Assistant",       "26-tool agent that reads, books, and updates everything."),
    ]

    for i, (icon, title, body) in enumerate(features):
        row = i // 2
        col = i % 2
        cx = start_x + col * (card_w + col_gap)
        cy = top_y + row * (card_h + row_gap)
        feature_card(slide, cx, cy, card_w, card_h, icon, title, body)

    # 4-step horizontal flow
    flow_y = top_y + 2 * card_h + 2 * row_gap + 182880
    flow_h = H - flow_y - 228600
    steps = ["\u2460 Profile", "\u2461 Benefits", "\u2462 Calendar", "\u2463 Done"]
    step_w = 1828800
    arrow_w = 365760
    total_flow_w = len(steps) * step_w + (len(steps) - 1) * arrow_w
    flow_start_x = (W - total_flow_w) // 2

    for i, step in enumerate(steps):
        sx = flow_start_x + i * (step_w + arrow_w)
        add_rect(slide, sx, flow_y, step_w, flow_h, SLATE1)
        tb = add_textbox(slide, sx, flow_y, step_w, flow_h)
        tf = tb.text_frame
        tf.word_wrap = False
        clear_tf(tf)
        set_para(tf, step, 13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        if i < len(steps) - 1:
            ax = sx + step_w + 45720
            ay = flow_y + flow_h // 2 - 91440
            tb_arr = add_textbox(slide, ax, ay, arrow_w - 91440, 182880)
            tf_arr = tb_arr.text_frame
            tf_arr.word_wrap = False
            clear_tf(tf_arr)
            set_para(tf_arr, "\u2192", 18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# SLIDE 4 — TECHNICAL + AI
# ---------------------------------------------------------------------------
def build_slide4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    section_pill(slide, "THE TECH")
    slide_title(slide, "Built to production standard in 24 hours")
    blue_underline(slide, w=5486400)

    card_top = 1828800
    card_h   = 4114800
    card_gap = 274320
    card_w   = (W - 1097280 - card_gap) // 2
    left_x   = 548640
    right_x  = left_x + card_w + card_gap

    STRIP = 54864

    # LEFT card — Stack
    add_rect(slide, left_x, card_top, card_w, card_h, SLATE1)
    add_rect(slide, left_x, card_top, STRIP, card_h, BLUE)

    tb_l = add_textbox(slide, left_x + STRIP + 137160, card_top + 182880,
                        card_w - STRIP - 228600, card_h - 365760)
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    clear_tf(tf_l)
    set_para(tf_l, "Stack", 16, bold=True, color=WHITE)
    stack_items = [
        "React 19 + Vite (frontend)",
        "FastAPI + Python (backend)",
        "Firebase Auth + Firestore",
        "Leaflet + Google Maps API",
        "Anthropic Claude Haiku",
        "Resend API (email confirmations)",
        "Google Calendar OAuth2",
        "pypdf (insurance PDF parsing)",
    ]
    for item in stack_items:
        set_para(tf_l, f"\u2022  {item}", 11, bold=False, color=SLATE2, space_before=3)

    # RIGHT card — AI Agent
    add_rect(slide, right_x, card_top, card_w, card_h, SLATE1)
    add_rect(slide, right_x, card_top, STRIP, card_h, BLUE)

    tb_r = add_textbox(slide, right_x + STRIP + 137160, card_top + 182880,
                        card_w - STRIP - 228600, card_h - 365760)
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    clear_tf(tf_r)
    set_para(tf_r, "AI Agent \u2014 26 Tools", 16, bold=True, color=WHITE)
    agent_items = [
        ("READ:", "profile, benefits, appointments, health history, clinics"),
        ("WRITE:", "book/cancel, update profile, allergies, checkup dates, calendar"),
        ("SYSTEM:", "navigate app, notifications, employer benefit codes"),
    ]
    for label, detail in agent_items:
        p = tf_r.add_paragraph()
        p.space_before = Pt(6)
        run_label = p.add_run()
        run_label.text = label + "  "
        run_label.font.size = Pt(11)
        run_label.font.bold = True
        run_label.font.color.rgb = BLUE
        run_detail = p.add_run()
        run_detail.text = detail
        run_detail.font.size = Pt(11)
        run_detail.font.bold = False
        run_detail.font.color.rgb = SLATE2

    # Bottom italic note
    note_y = card_top + card_h + 182880
    caption(slide,
            "All Anthropic calls are server-side \u2014 API key never reaches the browser",
            548640, note_y, W - 1097280, 320040,
            italic=True, align=PP_ALIGN.CENTER, size=12)


# ---------------------------------------------------------------------------
# SLIDE 5 — BUSINESS + CLOSE
# ---------------------------------------------------------------------------
def build_slide5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    section_pill(slide, "THE BUSINESS")
    slide_title(slide, "Two revenue tracks. Clear ROI.")
    blue_underline(slide, w=4114800)

    panel_top  = 1828800
    panel_h    = 2926080
    panel_gap  = 274320
    panel_w    = (W - 1097280 - panel_gap) // 2
    left_x     = 548640
    right_x    = left_x + panel_w + panel_gap
    STRIP      = 54864

    # LEFT panel — Individual (plain slate card)
    add_rect(slide, left_x, panel_top, panel_w, panel_h, SLATE1)
    add_rect(slide, left_x, panel_top, STRIP, panel_h, BLUE)

    tb_ind = add_textbox(slide, left_x + STRIP + 137160, panel_top + 182880,
                          panel_w - STRIP - 228600, panel_h - 365760)
    tf_ind = tb_ind.text_frame
    tf_ind.word_wrap = True
    clear_tf(tf_ind)
    set_para(tf_ind, "Individual", 16, bold=True, color=WHITE)
    ind_items = [
        "Free: scheduling, map, basic AI",
        "Premium $8/mo: full AI agent,",
        "PDF parsing, history",
    ]
    for item in ind_items:
        set_para(tf_ind, f"\u2022  {item}", 12, bold=False, color=SLATE2, space_before=5)

    # RIGHT panel — B2B Employer (blue-tint card, PRIMARY badge)
    # Slightly lighter fill to signal primary
    b2b_fill = RGBColor(0x0D, 0x1F, 0x3C)
    add_rect(slide, right_x, panel_top, panel_w, panel_h, b2b_fill)
    add_rect(slide, right_x, panel_top, STRIP, panel_h, BLUE)
    # Top blue border for highlight
    add_rect(slide, right_x, panel_top, panel_w, 54864, BLUE)

    # PRIMARY badge
    badge_w = 913320
    badge_h = 228600
    badge_x = right_x + panel_w - badge_w - 182880
    badge_y = panel_top + 137160
    add_rect(slide, badge_x, badge_y, badge_w, badge_h, BLUE)
    tb_badge = add_textbox(slide, badge_x, badge_y, badge_w, badge_h)
    tf_badge = tb_badge.text_frame
    tf_badge.word_wrap = False
    clear_tf(tf_badge)
    set_para(tf_badge, "PRIMARY", 9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    tb_b2b = add_textbox(slide, right_x + STRIP + 137160, panel_top + 182880,
                          panel_w - STRIP - 228600, panel_h - 365760)
    tf_b2b = tb_b2b.text_frame
    tf_b2b.word_wrap = True
    clear_tf(tf_b2b)
    set_para(tf_b2b, "B2B Employer", 16, bold=True, color=WHITE)
    b2b_items = [
        "Per-seat team onboarding",
        "HR benefit utilization dashboard",
        "67% of employees don't use their",
        "benefits \u2014 employers pay every year",
    ]
    for item in b2b_items:
        set_para(tf_b2b, f"\u2022  {item}", 12, bold=False, color=SLATE2, space_before=5)

    # Divider line
    div_y = panel_top + panel_h + 228600
    add_rect(slide, 548640, div_y, W - 1097280, 27432, SLATE1)

    # Phase 2 note
    phase_y = div_y + 91440
    tb_p2 = add_textbox(slide, 548640, phase_y, W - 1097280, 320040)
    tf_p2 = tb_p2.text_frame
    tf_p2.word_wrap = True
    clear_tf(tf_p2)
    p_phase = tf_p2.paragraphs[0]
    p_phase.alignment = PP_ALIGN.CENTER
    r_label = p_phase.add_run()
    r_label.text = "Phase 2 \u2192  "
    r_label.font.size = Pt(13)
    r_label.font.bold = True
    r_label.font.color.rgb = BLUE
    r_detail = p_phase.add_run()
    r_detail.text = "Clinic partnerships: live booking + referral revenue"
    r_detail.font.size = Pt(13)
    r_detail.font.bold = False
    r_detail.font.color.rgb = SLATE2

    # Final CTA
    cta_y = phase_y + 365760
    tb_cta = add_textbox(slide, 0, cta_y, W, 548640)
    tf_cta = tb_cta.text_frame
    tf_cta.word_wrap = False
    clear_tf(tf_cta)
    set_para(tf_cta, "Let us show you.", 36, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    # Footer
    footer_y = H - 320040
    tb_foot = add_textbox(slide, 0, footer_y, W, 274320)
    tf_foot = tb_foot.text_frame
    tf_foot.word_wrap = False
    clear_tf(tf_foot)
    set_para(tf_foot, "Hack the Globe 2026", 10, bold=False, color=SLATE2, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    prs = new_prs()
    build_slide1(prs)
    build_slide2(prs)
    build_slide3(prs)
    build_slide4(prs)
    build_slide5(prs)
    out = "/Users/nicol/Desktop/NexaCare/NexaCare_Pitch_5.pptx"
    prs.save(out)
    print(f"Saved: {out}")
