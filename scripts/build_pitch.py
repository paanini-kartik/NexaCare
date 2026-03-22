"""
NexaCare Pitch Deck Builder
Generates NexaCare_Pitch.pptx with 8 slides
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import math

# ── Colour palette ─────────────────────────────────────────────
NAVY      = RGBColor(0x0f, 0x17, 0x2a)   # background
BLUE      = RGBColor(0x3b, 0x82, 0xf6)   # accent
SLATE     = RGBColor(0x94, 0xa3, 0xb8)   # secondary text
WHITE     = RGBColor(0xff, 0xff, 0xff)   # primary text
DARK_CARD = RGBColor(0x1e, 0x29, 0x3b)   # card bg
BLUE_DARK = RGBColor(0x1d, 0x4e, 0xd8)   # darker blue for gradient simulation
GREEN     = RGBColor(0x10, 0xb9, 0x81)   # success
CARD_BDR  = RGBColor(0x2d, 0x3f, 0x5e)   # card border

# ── Slide dimensions (widescreen 16:9) ─────────────────────────
W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK_LAYOUT = prs.slide_layouts[6]   # completely blank


# ════════════════════════════════════════════════════════════════
# HELPERS
# ════════════════════════════════════════════════════════════════

def add_bg(slide, color=NAVY):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=Pt(0)):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=Pt(18), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name = "Calibri"
    return txBox


def add_multiline_text(slide, lines, x, y, w, h,
                       font_size=Pt(16), bold=False, color=WHITE,
                       align=PP_ALIGN.LEFT, line_spacing=None):
    """Add text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txBox


def accent_bar(slide, y=Inches(0.07), width=Inches(2.2)):
    """Blue accent bar at top-left."""
    add_rect(slide, Inches(0.5), y, width, Inches(0.045), BLUE)


def slide_label(slide, label):
    """Small ALL-CAPS section label top-right."""
    add_text(slide, label,
             W - Inches(2.8), Inches(0.18), Inches(2.6), Inches(0.35),
             font_size=Pt(9), color=SLATE, align=PP_ALIGN.RIGHT)


def bottom_rule(slide):
    """Thin blue line near bottom."""
    add_rect(slide, Inches(0.5), H - Inches(0.55), W - Inches(1.0), Inches(0.025), BLUE)


def add_card(slide, x, y, w, h, fill=DARK_CARD, border=BLUE, bw=Pt(1.2)):
    return add_rect(slide, x, y, w, h, fill, border, bw)


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — HOOK
# ════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s1)

# Subtle grid lines for depth
for i in range(1, 5):
    add_rect(s1, Inches(i * 2.8), 0, Inches(0.012), H, RGBColor(0x1a, 0x24, 0x38))
for i in range(1, 4):
    add_rect(s1, 0, Inches(i * 2.0), W, Inches(0.012), RGBColor(0x1a, 0x24, 0x38))

# Large "N" logomark — circle + letter
cx, cy = Inches(6.667), Inches(2.6)
r = Inches(1.05)
logo_circle = add_rect(s1, cx - r, cy - r, r * 2, r * 2, BLUE)
logo_circle.line.fill.background()
# Round the circle via shape adjustments is not trivially available;
# use an ellipse (MSO_SHAPE = 9)
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu as E

def add_ellipse(slide, x, y, w, h, fill_color, line_color=None, lw=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(9, x, y, w, h)  # 9 = oval
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = lw
    else:
        shape.line.fill.background()
    return shape

# Remove rect, draw ellipse instead
sp = s1.shapes._spTree
sp.remove(logo_circle._element)

add_ellipse(s1, cx - r, cy - r, r * 2, r * 2, BLUE)

# "N" letter inside
add_text(s1, "N",
         cx - r, cy - r, r * 2, r * 2,
         font_size=Pt(72), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Title
add_text(s1, "NexaCare",
         Inches(2.5), Inches(4.0), Inches(8.333), Inches(1.1),
         font_size=Pt(60), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Tagline
add_text(s1, '"Your benefits exist. Start using them."',
         Inches(2.5), Inches(5.1), Inches(8.333), Inches(0.75),
         font_size=Pt(22), bold=False, color=BLUE, align=PP_ALIGN.CENTER, italic=True)

# Bottom
add_rect(s1, Inches(3.5), H - Inches(0.95), Inches(6.333), Inches(0.025), SLATE)
add_text(s1, "Hack the Globe 2025  ·  Team NexaCare",
         Inches(2.5), H - Inches(0.88), Inches(8.333), Inches(0.45),
         font_size=Pt(12), color=SLATE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM (stats)
# ════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s2)
accent_bar(s2)
slide_label(s2, "THE PROBLEM")

add_text(s2, "The friction is the problem",
         Inches(0.5), Inches(0.55), Inches(10), Inches(0.85),
         font_size=Pt(40), bold=True, color=WHITE)

# 3 stat cards
stats = [
    ("42%", "of adults skipped a\npreventive checkup last year"),
    ("67%", "of employees don't\nfully use employer benefits"),
    ("$88B", "spent annually on\npreventable conditions"),
]

card_w = Inches(3.6)
card_h = Inches(3.2)
gap    = Inches(0.4)
start_x = (W - (card_w * 3 + gap * 2)) / 2
card_y = Inches(1.65)

for i, (num, label) in enumerate(stats):
    cx = start_x + i * (card_w + gap)
    add_card(s2, cx, card_y, card_w, card_h)
    # Blue top strip
    add_rect(s2, cx, card_y, card_w, Inches(0.08), BLUE)
    # Big number
    add_text(s2, num,
             cx, card_y + Inches(0.35), card_w, Inches(1.3),
             font_size=Pt(58), bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    # Label
    add_multiline_text(s2, label.split('\n'),
                       cx + Inches(0.2), card_y + Inches(1.7), card_w - Inches(0.4), Inches(1.3),
                       font_size=Pt(15), color=SLATE, align=PP_ALIGN.CENTER)

bottom_rule(s2)
add_text(s2,
         "The coverage exists. The clinics exist. The layer connecting them doesn't.",
         Inches(0.5), H - Inches(0.52), W - Inches(1.0), Inches(0.38),
         font_size=Pt(12), color=SLATE, align=PP_ALIGN.CENTER, italic=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — PROBLEM DEEP (chaos map)
# ════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s3)
accent_bar(s3)
slide_label(s3, "THE PROBLEM")

add_text(s3, 'What "having benefits" actually looks like',
         Inches(0.5), Inches(0.55), Inches(11), Inches(0.85),
         font_size=Pt(36), bold=True, color=WHITE)

# Helper for chaos map nodes
def chaos_node(slide, x, y, w, h, title, icon, border_color=CARD_BDR):
    add_card(slide, x, y, w, h, fill=DARK_CARD, border=border_color, bw=Pt(1.5))
    add_text(slide, icon, x, y + Inches(0.12), w, Inches(0.5),
             font_size=Pt(22), align=PP_ALIGN.CENTER, color=WHITE)
    add_text(slide, title, x, y + Inches(0.6), w, Inches(0.7),
             font_size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Portal nodes (left column)
node_w = Inches(2.3)
node_h = Inches(1.15)

portals = [
    (Inches(0.55), Inches(1.7),  "Dental Portal",   "🦷", BLUE),
    (Inches(0.55), Inches(3.1),  "Vision Portal",   "👁", BLUE),
    (Inches(0.55), Inches(4.5),  "Physio Portal",   "🏃", BLUE),
]

for px, py, title, icon, bc in portals:
    chaos_node(s3, px, py, node_w, node_h, title, icon, bc)

# Right side nodes
chaos_node(s3, Inches(9.8),  Inches(2.1),  node_w, node_h, "Benefits PDF",       "❓", RGBColor(0xf5, 0x9e, 0x0b))
chaos_node(s3, Inches(9.8),  Inches(4.0),  node_w, node_h, "Annual Reset Date",  "⏰", RGBColor(0xef, 0x44, 0x44))

# Center frustrated user
user_r = Inches(0.75)
ucx, ucy = Inches(6.5), Inches(3.35)
add_ellipse(s3, ucx - user_r, ucy - user_r, user_r * 2, user_r * 2,
            RGBColor(0x27, 0x38, 0x50), SLATE, Pt(1.5))
add_text(s3, "😤", ucx - user_r, ucy - user_r + Inches(0.1), user_r * 2, user_r * 2,
         font_size=Pt(28), align=PP_ALIGN.CENTER, color=WHITE)
add_text(s3, "You",
         ucx - Inches(0.5), ucy + user_r + Inches(0.05), Inches(1.0), Inches(0.35),
         font_size=Pt(11), color=SLATE, align=PP_ALIGN.CENTER)

# Connector lines (arrows simulated with thin rects + dots)
def h_connector(slide, x1, y1, x2, y2):
    """Draw a simple L-shaped connector."""
    mid_x = (x1 + x2) / 2
    # horizontal part from x1 to mid
    lw = Inches(0.018)
    if abs(y2 - y1) < Inches(0.05):
        add_rect(slide, min(x1, x2), y1 - lw/2, abs(x2 - x1), lw, SLATE)
    else:
        add_rect(slide, min(x1, mid_x), y1 - lw/2, abs(mid_x - x1), lw, SLATE)
        add_rect(slide, mid_x - lw/2, min(y1, y2), lw, abs(y2 - y1), SLATE)
        add_rect(slide, min(mid_x, x2), y2 - lw/2, abs(x2 - mid_x), lw, SLATE)

# Portals → center
for px, py, *_ in portals:
    h_connector(s3,
                px + node_w,            py + node_h / 2,
                ucx - user_r,           ucy)

# Right nodes → center
h_connector(s3, Inches(9.8), Inches(2.1) + node_h / 2, ucx + user_r, ucy - Inches(0.2))
h_connector(s3, Inches(9.8), Inches(4.0) + node_h / 2, ucx + user_r, ucy + Inches(0.2))

bottom_rule(s3)
add_text(s3,
         "Three portals. One document you never opened. One date that resets every year.",
         Inches(0.5), H - Inches(0.52), W - Inches(1.0), Inches(0.38),
         font_size=Pt(12), color=SLATE, align=PP_ALIGN.CENTER, italic=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — SOLUTION (2×2 feature grid)
# ════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s4)
accent_bar(s4)
slide_label(s4, "THE SOLUTION")

add_text(s4, "One place. Four features.",
         Inches(0.5), Inches(0.55), Inches(10), Inches(0.85),
         font_size=Pt(40), bold=True, color=WHITE)

features = [
    ("🗓", "Smart Scheduling",      "Know what's overdue.\nBook in seconds."),
    ("💰", "Benefits Tracker",      "See exactly what's covered\nand what's left."),
    ("🗺", "Health Compass",        "Map of nearby clinics filtered\nby your actual coverage."),
    ("🤖", "AI Health Assistant",   "Ask anything about your\nbenefits, history, and care."),
]

fc_w = Inches(5.6)
fc_h = Inches(2.45)
fc_gap = Inches(0.3)
fc_start_x = (W - (fc_w * 2 + fc_gap)) / 2
fc_start_y = Inches(1.65)

for i, (icon, title, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    fx = fc_start_x + col * (fc_w + fc_gap)
    fy = fc_start_y + row * (fc_h + fc_gap * 0.8)

    add_card(s4, fx, fy, fc_w, fc_h, fill=DARK_CARD, border=BLUE, bw=Pt(1.5))
    # Left blue accent strip
    add_rect(s4, fx, fy, Inches(0.07), fc_h, BLUE)

    # Icon
    add_text(s4, icon, fx + Inches(0.2), fy + Inches(0.25), Inches(0.8), Inches(0.7),
             font_size=Pt(26), align=PP_ALIGN.CENTER)
    # Title
    add_text(s4, title,
             fx + Inches(1.05), fy + Inches(0.28), fc_w - Inches(1.2), Inches(0.55),
             font_size=Pt(19), bold=True, color=WHITE)
    # Description
    add_multiline_text(s4, desc.split('\n'),
                       fx + Inches(1.05), fy + Inches(0.88), fc_w - Inches(1.2), Inches(1.3),
                       font_size=Pt(14), color=SLATE)


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — HOW IT WORKS (timeline)
# ════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s5)
accent_bar(s5)
slide_label(s5, "HOW IT WORKS")

add_text(s5, "Two minutes to set up. Works forever after.",
         Inches(0.5), Inches(0.55), Inches(11), Inches(0.85),
         font_size=Pt(38), bold=True, color=WHITE)

steps = [
    ("01", "Profile",   "Age, occupation,\nlast checkup dates", "30 sec"),
    ("02", "Benefits",  "Employer invite code\nor insurance PDF",  "30 sec"),
    ("03", "Calendar",  "Connect Google\nCalendar (optional)",    "10 sec"),
    ("04", "Done",      "Dashboard: overdue\nchecks, clinics, AI", "🎉"),
]

step_w = Inches(2.8)
step_h = Inches(3.4)
step_gap = Inches(0.4)
step_start_x = (W - (step_w * 4 + step_gap * 3)) / 2
step_y = Inches(1.6)

# Connector line behind cards
conn_y = step_y + Inches(0.52)
add_rect(s5,
         step_start_x + step_w / 2,
         conn_y,
         step_w * 3 + step_gap * 3,
         Inches(0.03), BLUE)

for i, (num, title, body, timing) in enumerate(steps):
    sx = step_start_x + i * (step_w + step_gap)

    is_done = (i == 3)
    card_fill = RGBColor(0x1d, 0x4e, 0xd8) if is_done else DARK_CARD
    border_c  = WHITE if is_done else BLUE

    add_card(s5, sx, step_y, step_w, step_h, fill=card_fill, border=border_c, bw=Pt(1.5))

    # Step number circle
    circle_r = Inches(0.38)
    circle_x = sx + step_w / 2 - circle_r
    add_ellipse(s5, circle_x, step_y + Inches(0.15), circle_r * 2, circle_r * 2,
                BLUE if not is_done else WHITE)
    add_text(s5, num,
             circle_x, step_y + Inches(0.15), circle_r * 2, circle_r * 2,
             font_size=Pt(13), bold=True,
             color=WHITE if not is_done else BLUE,
             align=PP_ALIGN.CENTER)

    # Title
    add_text(s5, title,
             sx + Inches(0.15), step_y + Inches(1.05), step_w - Inches(0.3), Inches(0.55),
             font_size=Pt(18), bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)

    # Body
    add_multiline_text(s5, body.split('\n'),
                       sx + Inches(0.15), step_y + Inches(1.65), step_w - Inches(0.3), Inches(1.1),
                       font_size=Pt(13), color=SLATE if not is_done else WHITE,
                       align=PP_ALIGN.CENTER)

    # Timing badge
    add_text(s5, timing,
             sx + Inches(0.15), step_y + step_h - Inches(0.55), step_w - Inches(0.3), Inches(0.4),
             font_size=Pt(12), color=BLUE if not is_done else WHITE,
             bold=True, align=PP_ALIGN.CENTER)

bottom_rule(s5)
add_text(s5,
         "After setup: every appointment booked syncs to your calendar. Confirmation email sent automatically.",
         Inches(0.5), H - Inches(0.52), W - Inches(1.0), Inches(0.38),
         font_size=Pt(11.5), color=SLATE, align=PP_ALIGN.CENTER, italic=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — AI LAYER
# ════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s6)
accent_bar(s6)
slide_label(s6, "AI LAYER")

add_text(s6, "An assistant that actually knows you",
         Inches(0.5), Inches(0.55), Inches(11), Inches(0.85),
         font_size=Pt(38), bold=True, color=WHITE)

col_w = Inches(5.5)
col_h = Inches(4.35)
col_y = Inches(1.55)
gap   = Inches(0.45)
left_x  = (W - col_w * 2 - gap) / 2
right_x = left_x + col_w + gap

# Left card — "What the AI knows"
add_card(s6, left_x, col_y, col_w, col_h)
add_rect(s6, left_x, col_y, col_w, Inches(0.07), BLUE)
add_text(s6, "What the AI knows",
         left_x + Inches(0.25), col_y + Inches(0.18), col_w - Inches(0.5), Inches(0.55),
         font_size=Pt(17), bold=True, color=BLUE)

knows = [
    "Your name, age, occupation, location",
    "Benefits used vs. remaining (dental, vision, physio)",
    "Upcoming and past appointments",
    "Allergies, medical history, last checkup dates",
]
for j, item in enumerate(knows):
    add_text(s6, f"• {item}",
             left_x + Inches(0.3), col_y + Inches(0.9) + j * Inches(0.75),
             col_w - Inches(0.5), Inches(0.65),
             font_size=Pt(13.5), color=WHITE)

# Right card — "What it can do"
add_card(s6, right_x, col_y, col_w, col_h)
add_rect(s6, right_x, col_y, col_w, Inches(0.07), BLUE)
add_text(s6, "What it can do",
         right_x + Inches(0.25), col_y + Inches(0.18), col_w - Inches(0.5), Inches(0.55),
         font_size=Pt(17), bold=True, color=BLUE)

cando = [
    "Book or cancel appointments",
    "Find nearby clinics by type",
    "Update your health profile by chat",
    "Read uploaded insurance PDFs",
    "Add appointments to Google Calendar",
    "Surface checkups you're overdue for",
]
for j, item in enumerate(cando):
    add_text(s6, f"• {item}",
             right_x + Inches(0.3), col_y + Inches(0.9) + j * Inches(0.57),
             col_w - Inches(0.5), Inches(0.52),
             font_size=Pt(13.5), color=WHITE)

bottom_rule(s6)
add_text(s6,
         "Powered by Claude (Anthropic). Fully server-side — your data never leaves NexaCare's backend.",
         Inches(0.5), H - Inches(0.52), W - Inches(1.0), Inches(0.38),
         font_size=Pt(11.5), color=SLATE, align=PP_ALIGN.CENTER, italic=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — BUSINESS MODEL
# ════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s7)
accent_bar(s7)
slide_label(s7, "BUSINESS MODEL")

add_text(s7, "Two revenue tracks",
         Inches(0.5), Inches(0.55), Inches(10), Inches(0.85),
         font_size=Pt(40), bold=True, color=WHITE)

panel_y = Inches(1.55)
panel_h = Inches(4.0)
panel_w = Inches(5.6)
gap_p   = Inches(0.45)
left_px  = (W - panel_w * 2 - gap_p) / 2
right_px = left_px + panel_w + gap_p

# Individual panel
add_card(s7, left_px, panel_y, panel_w, panel_h, border=CARD_BDR, bw=Pt(1))
add_rect(s7, left_px, panel_y, panel_w, Inches(0.07), SLATE)
add_text(s7, "Individual",
         left_px + Inches(0.25), panel_y + Inches(0.18), panel_w - Inches(0.5), Inches(0.52),
         font_size=Pt(17), bold=True, color=WHITE)

ind_items = [
    ("Free tier",   "Smart scheduling, Health Compass, basic AI"),
    ("Premium",     "$8 / month — full AI layer, PDF parsing,\nunlimited history"),
]
for j, (badge, desc) in enumerate(ind_items):
    by = panel_y + Inches(0.88) + j * Inches(1.4)
    add_rect(s7, left_px + Inches(0.25), by, Inches(1.1), Inches(0.36),
             DARK_CARD, SLATE, Pt(0.8))
    add_text(s7, badge,
             left_px + Inches(0.25), by, Inches(1.1), Inches(0.36),
             font_size=Pt(11), bold=True, color=SLATE, align=PP_ALIGN.CENTER)
    add_multiline_text(s7, desc.split('\n'),
                       left_px + Inches(1.5), by, panel_w - Inches(1.7), Inches(0.9),
                       font_size=Pt(13), color=WHITE)

# B2B panel (highlighted)
add_card(s7, right_px, panel_y, panel_w, panel_h, fill=RGBColor(0x1a, 0x34, 0x5c), border=BLUE, bw=Pt(2.0))
add_rect(s7, right_px, panel_y, panel_w, Inches(0.07), BLUE)
# "Primary" badge
add_rect(s7, right_px + panel_w - Inches(1.3), panel_y - Inches(0.22),
         Inches(1.2), Inches(0.32), BLUE)
add_text(s7, "PRIMARY",
         right_px + panel_w - Inches(1.3), panel_y - Inches(0.22),
         Inches(1.2), Inches(0.32),
         font_size=Pt(9), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s7, "B2B — Employer Plans",
         right_px + Inches(0.25), panel_y + Inches(0.18), panel_w - Inches(0.5), Inches(0.52),
         font_size=Pt(17), bold=True, color=BLUE)

b2b_items = [
    "Team onboarding dashboard",
    "HR benefit utilization analytics",
    "Employee engagement reporting",
]
for j, item in enumerate(b2b_items):
    add_text(s7, f"• {item}",
             right_px + Inches(0.3), panel_y + Inches(0.88) + j * Inches(0.65),
             panel_w - Inches(0.5), Inches(0.58),
             font_size=Pt(14), color=WHITE)

add_text(s7,
         '"67% of employees don\'t use their benefits.\nEmployers pay for unused coverage every year."',
         right_px + Inches(0.25), panel_y + Inches(2.7), panel_w - Inches(0.5), Inches(1.1),
         font_size=Pt(12.5), color=SLATE, italic=True)

# Phase 2 arrow bar
p2_y = panel_y + panel_h + Inches(0.2)
add_rect(s7, left_px, p2_y, panel_w * 2 + gap_p, Inches(0.48), RGBColor(0x1e, 0x29, 0x3b),
         BLUE, Pt(0.8))
add_text(s7,
         "Phase 2 →  Clinic partnerships: live booking, real-time availability, referral revenue",
         left_px + Inches(0.2), p2_y + Inches(0.05), panel_w * 2 + gap_p - Inches(0.4), Inches(0.4),
         font_size=Pt(12.5), color=BLUE, bold=False)


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — CLOSE
# ════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s8)

# Grid lines for texture
for i in range(1, 5):
    add_rect(s8, Inches(i * 2.8), 0, Inches(0.012), H, RGBColor(0x1a, 0x24, 0x38))
for i in range(1, 4):
    add_rect(s8, 0, Inches(i * 2.0), W, Inches(0.012), RGBColor(0x1a, 0x24, 0x38))

# Large title
add_text(s8, "Built in 24 hours.",
         Inches(0.5), Inches(1.35), W - Inches(1.0), Inches(1.5),
         font_size=Pt(64), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Tech stack
add_text(s8, "Python  ·  React  ·  Firebase  ·  Google Calendar  ·  Anthropic Claude",
         Inches(1.0), Inches(2.95), W - Inches(2.0), Inches(0.6),
         font_size=Pt(16), color=SLATE, align=PP_ALIGN.CENTER)

# Role badges
roles  = ["Frontend", "Backend", "Maps", "AI"]
badge_w = Inches(2.2)
badge_h = Inches(0.52)
badge_gap = Inches(0.22)
badge_total_w = badge_w * 4 + badge_gap * 3
badge_start_x = (W - badge_total_w) / 2
badge_y = Inches(3.85)

for i, role in enumerate(roles):
    bx = badge_start_x + i * (badge_w + badge_gap)
    add_rect(s8, bx, badge_y, badge_w, badge_h, DARK_CARD, BLUE, Pt(1.2))
    add_text(s8, role, bx, badge_y, badge_w, badge_h,
             font_size=Pt(14), bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# Big CTA
add_text(s8, "Let us show you.",
         Inches(1.0), Inches(4.75), W - Inches(2.0), Inches(1.1),
         font_size=Pt(50), bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# Bottom
add_rect(s8, Inches(3.5), H - Inches(0.95), Inches(6.333), Inches(0.025), SLATE)
add_text(s8, "Hack the Globe 2025",
         Inches(1.0), H - Inches(0.88), W - Inches(2.0), Inches(0.45),
         font_size=Pt(12), color=SLATE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════
out_path = "/Users/nicol/Desktop/NexaCare/NexaCare_Pitch.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
